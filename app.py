import os
from io import BytesIO

import faiss
import httpx
import nlp
import numpy as np
import pandas as pd
import torch
from fuzzywuzzy import fuzz
from pptx import Presentation
from starlette.applications import Starlette
from starlette.config import environ
from starlette.middleware.cors import CORSMiddleware
from starlette.responses import JSONResponse, StreamingResponse

from lfqa_utils import (embed_passages_for_retrieval, make_qa_retriever_model,
                        make_qa_s2s_model, qa_s2s_generate,
                        query_mix_dense_index)
from tei import TEIFile

app = Starlette()

os.environ['KMP_DUPLICATE_LIB_OK'] = 'True'

device = 'cuda:0' if torch.cuda.is_available() else 'cpu'
print('Device:' + device)

if device == 'cuda:0':
    faiss_res = faiss.StandardGpuResources()

qar_tokenizer, qar_model = make_qa_retriever_model(
    model_name="google/bert_uncased_L-8_H-768_A-12",
    from_file='models/unfilter_ret_model_2.pth',
    device=device
)
qa_s2s_tokenizer, qa_s2s_model = make_qa_s2s_model(
    model_name="facebook/bart-large-cnn",
    from_file='models/keyword_cnn_5e-5_3_resave.pth',
    device=device
)

def build_snippet(tei):
    # Use a breakpoint in the code line below to debug your script.
    paper_snippet = {
        'article_title': [],
        'section_title': [],
        'passage_text': []}
    paper_snippet['article_title'].append('')
    paper_snippet['section_title'].append('Abstract')
    paper_snippet['passage_text'].append(tei.abstract)

    for h in tei.headers:
        sentences = tei.text[h['start']:h['end'] + 1]
        # four sentences averages close to 100 words
        fours = zip(*(iter(sentences),) * 4)
        for j in fours:
            tmp = ' '.join([x['string'] for x in j])
            paper_snippet['article_title'].append('')
            paper_snippet['section_title'].append(h['section'])
            paper_snippet['passage_text'].append(tmp)

    return paper_snippet


def build_indices(snippet, embedder, tokenizer):
    snippet = pd.DataFrame.from_dict(snippet)
    snippet = nlp.Dataset.from_pandas(snippet)

    passages = [p for p in snippet['passage_text']]
    keywords = [k for k in snippet['section_title']]

    passage_emb = embed_passages_for_retrieval(passages, tokenizer, embedder, 128, device)
    keyword_emb = embed_passages_for_retrieval(keywords, tokenizer, embedder, 128, device)

    return passage_emb, keyword_emb


def build_input(question, tei, snippet, passage_emb, keyword_emb):
    snippet = pd.DataFrame.from_dict(snippet)
    snippet = nlp.Dataset.from_pandas(snippet)
    passage = faiss.IndexFlatIP(128)
    keyword = faiss.IndexFlatIP(128)

    if device == 'cuda:0':
        passage = faiss.index_cpu_to_gpu(faiss_res, 0, passage)
        keyword = faiss.index_cpu_to_gpu(faiss_res, 0, keyword)

    passage.add(np.array(passage_emb, dtype=np.float32))
    keyword.add(np.array(keyword_emb, dtype=np.float32))
    doc, res_list = query_mix_dense_index(
        question, qar_model, qar_tokenizer,
        snippet, passage, keyword, weight=0.75, device=device
    )
    headers = [(elem['section'], elem['n']) for elem in tei.headers]
    scrs = [fuzz.ratio(question.lower(), elem[0].lower()) for elem in headers]
    keywords = []
    if len(scrs) > 0 and max(scrs) >= 90:
        begin = headers[scrs.index(max(scrs))]
        for elem in headers:
            if elem == begin:
                continue
            if begin[1] == elem[1][:len(begin[1])]:
                keywords.append(elem['section'])
        if len(keywords) > 0:
            question = question + ', ' + ', '.join(keywords)
    question_doc = "question: {} context: {}".format(question, doc)
    return question_doc


if 'DEBUG' in environ:
    app.add_middleware(CORSMiddleware, allow_origin_regex='http://localhost:.*',
                       allow_methods=['*'], allow_headers=['*'])
else:
    app.add_middleware(CORSMiddleware, allow_origin_regex='https://doc2slides.s3-web.us-east.cloud-object-storage.appdomain.cloud',
                       allow_methods=['*'], allow_headers=['*'])

@app.on_event("startup")
async def startup():
    pass


@app.on_event("shutdown")
async def shutdown():
    pass


@app.route('/encode_paper', methods=["POST"])
async def encode_paper(request):
    form = await request.form()
    content = await form['file'].read()
    async with httpx.AsyncClient() as client:
        r = await client.post('http://localhost:8070/api/processFulltextDocument',
                              files={'input': content}, timeout=None)
        tei = TEIFile(r.text)
        paper_snippet = build_snippet(tei)
        passage_emb, keyword_emb = build_indices(
            paper_snippet, qar_model, qar_tokenizer)
        
        return JSONResponse({
            'tei': r.text,
            'snippets': paper_snippet,
            'passage_embedding': passage_emb.tolist(),
            'keyword_embedding': keyword_emb.tolist()
        })


@app.route('/generate_slide', methods=["POST"])
async def gen_slide(request):
    req = await request.json()
    question_doc = build_input(req['title'], TEIFile(req['tei']), req['snippets'], req['passage_embedding'], req['keyword_embedding'])

    answer = qa_s2s_generate(
        question_doc, qa_s2s_model, qa_s2s_tokenizer,
        num_answers=1,
        num_beams=8,
        min_len=64,
        max_len=200,
        max_input_length=1024,
        device=device)[0]

    return JSONResponse({
        'answer': answer.split('\n')
    })


@app.route('/download_slides', methods=["POST"])
async def download_slides(request):
    req = await request.json()
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]
    output = BytesIO()
    for slideContent in req['slides']:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slideContent['title']

        tf = body_shape.text_frame

        for bullet in slideContent['content']:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1
        
        tf.fit_text(font_family="Arial")

    prs.save(output)
    output.seek(0)

    return StreamingResponse(output, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
